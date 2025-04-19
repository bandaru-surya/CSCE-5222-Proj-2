import pandas as pd, matplotlib.pyplot as plt, joblib, os, sys
from pathlib import Path
from pptx import Presentation

# -------- paths ---------------------------------------------------------
BASE        = Path(__file__).resolve().parent.parent
DATA_DIR    = BASE / "data"
MODEL_DIR   = BASE / "models"
RESULTS_DIR = BASE / "results"
PPTX_IN     = BASE / "presentation" / "Group26.pptx"
PPTX_OUT    = PPTX_IN.with_stem(PPTX_IN.stem + "_filled")

for d in (DATA_DIR, MODEL_DIR, RESULTS_DIR):
    d.mkdir(parents=True, exist_ok=True)

# -------- import project modules ---------------------------------------
sys.path.append(str((BASE / "src").resolve()))
from synthetic_data import generate_synthetic_patients
from train_models   import train_all_models

# -------- Step 1: data --------------------------------------------------
df = generate_synthetic_patients(500)                         # original
df.to_csv(DATA_DIR / "synthetic_data.csv", index=False)

# TODO: call your VAE here → df_aug
df_aug = df.copy()                                            # stub

# -------- Step 2: train models -----------------------------------------
perf_orig = train_all_models(df)      # saves pickles into MODEL_DIR
perf_aug  = train_all_models(df_aug)
perf = perf_orig.join(perf_aug,
                      lsuffix="_orig",
                      rsuffix="_vae")
perf.to_csv(RESULTS_DIR / "performance.csv")

# -------- Step 3: bar‑chart --------------------------------------------
plt.figure(figsize=(8,4))
perf["MAE_orig"].plot(kind="bar", label="Orig", alpha=.7)
perf["MAE_vae"].plot(kind="bar", label="VAE", alpha=.7)
plt.ylabel("MAE  (months)")
plt.title("Model comparison")
plt.legend()
plt.tight_layout()
plt.savefig(RESULTS_DIR / "performance_bar.png")
plt.close()

# -------- Step 4: update PPTX -------------------------------------------
MAE_best = perf["MAE_vae"].min()
improve  = 100 * (perf["MAE_orig"].min() - MAE_best) / perf["MAE_orig"].min()

replace_map = {
    "«Nreal»"       : str(len(df)),
    "«Nsyn»"        : str(len(df_aug) - len(df)),
    "«rows × cols»" : f"{df.shape[0]} × {df.shape[1]}",
    "«k»"           : "0",
    "«x»"           : f"{improve:.1f}",
    "«y»"           : f"{(perf.loc['ElasticNet','MAE_orig'] - perf.loc['ElasticNet','MAE_vae']):.2f}",
    "«Δ»"           : f"{improve:.1f}",
    "«overall %»"   : f"{improve:.1f}",
}

prs = Presentation(PPTX_IN)
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        for key, val in replace_map.items():
            text = text.replace(key, val)
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.text = text

prs.save(PPTX_OUT)
print("✓ Pipeline finished — see:", RESULTS_DIR, "and", PPTX_OUT.name)